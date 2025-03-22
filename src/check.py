from pptx import Presentation
from google import genai
from google.genai import types
from src.utils import logger

class Checker:
    def __init__(self, api_key: str):
        self.client = genai.Client(api_key=api_key)
        self.logger = logger.setup_logger()

    def import_pptx(file_path):
        """
        PowerPointファイル（.pptx）を取り込む関数
        
        Args:
            file_path (str): 取り込むPowerPointファイルのパス
            
        Returns:
            Presentation: python-pptxのPresentationオブジェクト
        """
        
        try:
            presentation = Presentation(file_path)
            return presentation
        except Exception as e:
            print(f"ファイルの取り込み中にエラーが発生しました: {e}")
            return None


    def extract_pptx(self, temp_path: str):
        """
        PowerPointファイル（.pptx）からテキスト、フォント、フォントサイズ、フォントカラー、テキストボックスの位置を抽出する関数
        
        Args:
            temp_path: 取り込むPowerPointファイルのパス
            
        Returns:
            list: 各スライドの情報を含む辞書のリスト
        """
        # PowerPointファイルの取り込み
        try:
            presentation = Presentation(temp_path)
        except Exception as e:
            self.logger.error(f"PowerPointファイルの取り込み中にエラーが発生: {str(e)}", exc_info=True)
            raise
        
        slides_data = []
        
        for slide_index, slide in enumerate(presentation.slides):
            slide_data = {
                'slide_number': slide_index + 1,
                'shapes': []
            }
            
            for shape in slide.shapes:
                # テキストを含む形状のみを処理
                if not shape.has_text_frame:
                    continue
                    
                shape_data = {
                    'text': '',
                    'paragraphs': [],
                    'position': {
                        'left': shape.left,
                        'top': shape.top,
                        'width': shape.width,
                        'height': shape.height
                    }
                }
                
                for paragraph in shape.text_frame.paragraphs:
                    paragraph_text = paragraph.text
                    shape_data['text'] += paragraph_text + '\n'
                    
                    paragraph_info = {
                        'text': paragraph_text,
                        'runs': []
                    }
                    
                    for run in paragraph.runs:
                        font = run.font
                        color = 'なし'
                        
                        # フォントカラーの取得
                        if font.color.type is not None:
                            if hasattr(font.color, 'rgb') and font.color.rgb is not None:
                                rgb = font.color.rgb
                                color = f'RGB({rgb[0]}, {rgb[1]}, {rgb[2]})'
                        
                        run_info = {
                            'text': run.text,
                            'font_name': font.name,
                            'font_size': font.size.pt if font.size is not None else 'デフォルト',
                            'font_color': color,
                            'bold': font.bold,
                            'italic': font.italic,
                            'underline': font.underline
                        }
                        
                        paragraph_info['runs'].append(run_info)
                    
                    shape_data['paragraphs'].append(paragraph_info)
                
                slide_data['shapes'].append(shape_data)
            
            slides_data.append(slide_data)
        
        return slides_data

    def print_pptx(self, slides_data: list):
        """
        抽出したPowerPointの内容を表示する関数
        
        Args:
            slides_data : extract_pptx関数で抽出したデータ
        """
        
        if not slides_data:
            print("データがありません。")
            return
        
        for slide in slides_data:
            print(f"\n===== スライド {slide['slide_number']} =====")
            
            for shape_index, shape in enumerate(slide['shapes']):
                print(f"\n-- テキストボックス {shape_index + 1} --")
                print(f"位置: 左={shape['position']['left']}, 上={shape['position']['top']}, "
                        f"幅={shape['position']['width']}, 高さ={shape['position']['height']}")
                
                for para_index, paragraph in enumerate(shape['paragraphs']):
                    print(f"\n段落 {para_index + 1}:")
                    
                    for run_index, run in enumerate(paragraph['runs']):
                        print(f"  テキスト: {run['text']}")
                        print(f"  フォント: {run['font_name']}, サイズ: {run['font_size']}, 色: {run['font_color']}")
                        
                        style = []
                        if run['bold']:
                            style.append('太字')
                        if run['italic']:
                            style.append('斜体')
                        if run['underline']:
                            style.append('下線')
                        
                        if style:
                            print(f"  スタイル: {', '.join(style)}")
                        print("")


    def check_pptx(self, model: str, content: list, prompt: str) -> str:
        """PowerPointの内容をLLMでチェックする

        Args:
            model : 使用するモデル
            content : PowerPointの内容（スライドごとのテキストのリスト）
            prompt : チェックのためのプロンプト

        Returns:
            str: LLMからの分析結果
        """
        try:
            # スライドの内容を文字列に変換
            slides_text = "\n\n".join([
                f"スライド {i+1}:\n" + "\n".join(slide)
                for i, slide in enumerate(content)
            ])
            
            # プロンプトの作成
            full_prompt = f"""
    以下のPowerPointの内容を分析し、{prompt}の観点から評価してください。

    PowerPointの内容:
    {slides_text}

    分析結果を日本語で出力してください。
    """
            self.logger.info(full_prompt)
            # LLMによる分析の実行
            response = self.client.models.generate_content(
                model=model,
                contents=full_prompt
                )
            result = response.text
            
            self.logger.info("PowerPointの内容の分析が完了しました")
            return result
            
        except Exception as e:
            self.logger.error(f"LLMによる分析中にエラーが発生: {str(e)}", exc_info=True)
            raise
